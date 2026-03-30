"""
Assemble 10 PNG slides into a PowerPoint presentation.
"""
from glob import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


OUTPUT_DIR = Path("/Users/jongchan/Desktop/claude/in-house-clustering/output")
SLIDES_DIR = OUTPUT_DIR / "slides"
PPTX_PATH = OUTPUT_DIR / "clustering_methods_presentation.pptx"

BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_CC = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_88 = RGBColor(0x88, 0x88, 0x88)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG_COLOR)

    # Title text box — centered, upper portion
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11.333), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "In-House Layout Feature Clustering"
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(1), Inches(3.6), Inches(11.333), Inches(0.8)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "10가지 클러스터링 방법론 비교 분석"
    run2.font.size = Pt(24)
    run2.font.color.rgb = GRAY_CC

    # Sub-subtitle
    txBox3 = slide.shapes.add_textbox(
        Inches(1), Inches(4.55), Inches(11.333), Inches(0.7)
    )
    tf3 = txBox3.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "CD (nm) 산포 최소화 | SHAP Feature Engineering | Cost Function 최적화"
    run3.font.size = Pt(16)
    run3.font.color.rgb = GRAY_88


def add_image_slide(prs: Presentation, img_path: Path):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG_COLOR)

    slide.shapes.add_picture(
        str(img_path),
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    add_title_slide(prs)

    # Image slides sorted by name (slide_01 ... slide_10)
    png_files = sorted(glob(str(SLIDES_DIR / "slide_*.png")))
    if len(png_files) != 10:
        raise RuntimeError(f"Expected 10 PNGs, found {len(png_files)}: {png_files}")

    for png in png_files:
        add_image_slide(prs, Path(png))

    prs.save(str(PPTX_PATH))
    print(f"Saved: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
